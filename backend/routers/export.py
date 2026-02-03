"""PDF and Excel export endpoints for the Banking ML Platform."""

import io
import json
import logging
import tempfile
import zipfile
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional

from fastapi import APIRouter, Depends
from fastapi.responses import StreamingResponse
from pydantic import BaseModel

from backend.core.config import Settings
from backend.core.dependencies import get_settings
from backend.core.exceptions import NotFoundError, DataError, ValidationError

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/api/admin/export", tags=["export"])

# ── Helpers ──────────────────────────────────────────────────────────────────

def _find_uc_dir(uc_id: str, preprocess_dir: Path) -> Optional[Path]:
    """Resolve a use-case ID (UC-FR-01, uc_fr_01_fraud_scoring, etc.) to its
    preprocessing_output directory.  Tries exact match first, then prefix match."""
    if not preprocess_dir.exists():
        return None

    # Normalise: UC-FR-01 -> uc_fr_01
    normalised = uc_id.lower().replace("-", "_")

    for d in sorted(preprocess_dir.iterdir()):
        if not d.is_dir():
            continue
        dir_lower = d.name.lower()
        # exact match
        if dir_lower == normalised:
            return d
        # prefix match  (uc_fr_01 matches uc_fr_01_fraud_scoring)
        if dir_lower.startswith(normalised):
            return d
    return None


def _load_json(path: Path) -> Any:
    """Load a JSON file, return None on failure."""
    if not path.exists():
        return None
    try:
        return json.loads(path.read_text())
    except Exception as e:
        logger.warning("Failed to load %s: %s", path, e)
        return None


def _safe(val: Any, default: str = "N/A") -> str:
    """Convert a value to a display string, handling None gracefully."""
    if val is None:
        return default
    if isinstance(val, float):
        return f"{val:.2f}"
    return str(val)


def _trunc(text: str, max_len: int = 90) -> str:
    """Truncate text for table cells."""
    if len(text) <= max_len:
        return text
    return text[: max_len - 3] + "..."


# ── PDF generation ───────────────────────────────────────────────────────────

def _build_use_case_pdf(uc_dir: Path) -> bytes:
    """Generate a professional PDF report for a single use case."""
    from fpdf import FPDF

    summary = _load_json(uc_dir / "summary.json") or {}
    full_report = _load_json(uc_dir / "full_report.json") or {}
    column_profiles = _load_json(uc_dir / "column_profiles.json") or []
    feature_eng = _load_json(uc_dir / "feature_engineering.json") or []
    outliers = _load_json(uc_dir / "outliers.json") or []
    target_dist = _load_json(uc_dir / "target_distribution.json") or {}
    correlations = _load_json(uc_dir / "correlations.json") or []
    training_results = _load_json(uc_dir / "training_results.json")

    # Use full_report's data_quality dict if present
    dq = full_report.get("data_quality", {})

    label = summary.get("label", uc_dir.name)
    uc_key = summary.get("use_case_key", uc_dir.name)

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=20)

    # ── Colours ──
    HEADER_BG = (15, 52, 96)       # dark navy
    HEADER_FG = (255, 255, 255)
    SECTION_BG = (230, 240, 250)
    TABLE_HEADER_BG = (41, 82, 130)
    TABLE_HEADER_FG = (255, 255, 255)
    ROW_ALT = (245, 248, 252)
    ACCENT = (22, 119, 255)

    def add_header_bar():
        """Top header band on the current page."""
        pdf.set_fill_color(*HEADER_BG)
        pdf.rect(0, 0, 210, 28, "F")
        pdf.set_xy(10, 6)
        pdf.set_font("Helvetica", "B", 14)
        pdf.set_text_color(*HEADER_FG)
        pdf.cell(0, 8, "Banking AI/ML Platform  -  Use Case Report", ln=True)
        pdf.set_font("Helvetica", "", 9)
        pdf.set_x(10)
        pdf.cell(0, 5, f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}", ln=True)
        pdf.set_text_color(0, 0, 0)
        pdf.set_y(34)

    def section_title(title: str):
        """Draw a section heading."""
        pdf.ln(4)
        pdf.set_fill_color(*SECTION_BG)
        pdf.set_font("Helvetica", "B", 12)
        pdf.set_text_color(*HEADER_BG)
        pdf.cell(0, 9, f"  {title}", ln=True, fill=True)
        pdf.set_text_color(0, 0, 0)
        pdf.ln(2)

    def kv_row(key: str, value: str):
        """Key-value pair row."""
        pdf.set_font("Helvetica", "B", 9)
        pdf.cell(55, 6, key, border=0)
        pdf.set_font("Helvetica", "", 9)
        pdf.cell(0, 6, value, ln=True)

    def table_header(cols: List[tuple]):
        """Draw table header row. cols = [(label, width), ...]."""
        pdf.set_fill_color(*TABLE_HEADER_BG)
        pdf.set_text_color(*TABLE_HEADER_FG)
        pdf.set_font("Helvetica", "B", 8)
        for label, w in cols:
            pdf.cell(w, 7, label, border=0, fill=True)
        pdf.ln()
        pdf.set_text_color(0, 0, 0)

    def table_row(values: List[str], widths: List[int], alt: bool = False):
        """Draw a data row."""
        if alt:
            pdf.set_fill_color(*ROW_ALT)
        else:
            pdf.set_fill_color(255, 255, 255)
        pdf.set_font("Helvetica", "", 8)
        for val, w in zip(values, widths):
            pdf.cell(w, 6, _trunc(val, max_len=int(w * 1.3)), border=0, fill=True)
        pdf.ln()

    # ── Page 1 ────────────────────────────────────────────────────────────
    pdf.add_page()
    add_header_bar()

    # Use case identity
    section_title("Use Case Overview")
    kv_row("Name:", label)
    kv_row("Key:", uc_key)
    kv_row("Category:", summary.get("category", "N/A"))
    kv_row("Domain:", summary.get("domain", "N/A"))
    kv_row("Run Timestamp:", summary.get("run_timestamp", "N/A"))
    pdf.ln(2)

    # Data Quality
    section_title("Data Quality")
    kv_row("Total Rows:", _safe(summary.get("total_rows")))
    kv_row("Total Columns:", _safe(summary.get("total_columns")))
    kv_row("Numeric Columns:", _safe(summary.get("numeric_columns")))
    kv_row("Categorical Columns:", _safe(summary.get("categorical_columns")))
    kv_row("Quality Score:", f"{_safe(summary.get('data_quality_score'))} / 100")
    kv_row("Avg Missing %:", _safe(summary.get("avg_missing_pct")))
    kv_row("Duplicate Row %:", _safe(dq.get("duplicate_row_pct", summary.get("duplicate_row_pct"))))
    kv_row("Avg Outlier %:", _safe(summary.get("avg_outlier_pct")))

    missing_details = dq.get("missing_details", [])
    if missing_details:
        pdf.ln(1)
        pdf.set_font("Helvetica", "I", 8)
        pdf.cell(0, 5, f"Columns with missing data: {len(missing_details)}", ln=True)

    pdf.ln(2)

    # Target analysis
    section_title("Target Analysis")
    kv_row("Target Column:", _safe(target_dist.get("target_column", summary.get("target_column"))))
    kv_row("Number of Classes:", _safe(target_dist.get("n_classes")))
    kv_row("Imbalance Ratio:", _safe(target_dist.get("imbalance_ratio", summary.get("class_imbalance_ratio"))))

    class_dist = target_dist.get("class_distribution", {})
    if class_dist:
        pdf.ln(1)
        cols = [("Class", 50), ("Count", 40), ("Proportion", 40)]
        table_header(cols)
        total = sum(class_dist.values()) or 1
        widths = [c[1] for c in cols]
        for i, (cls, cnt) in enumerate(class_dist.items()):
            table_row([str(cls), str(cnt), f"{cnt / total * 100:.1f}%"], widths, alt=i % 2 == 1)

    pdf.ln(2)

    # Feature summary
    section_title("Feature Summary")
    numeric_profiles = [p for p in column_profiles if p.get("inferred_type") == "numeric"]
    cat_profiles = [p for p in column_profiles if p.get("inferred_type") == "categorical"]

    if numeric_profiles:
        pdf.set_font("Helvetica", "B", 9)
        pdf.cell(0, 6, f"Numeric Features ({len(numeric_profiles)})", ln=True)
        cols = [("Column", 40), ("Mean", 22), ("Std", 22), ("Min", 22), ("Median", 22), ("Max", 22), ("Missing%", 22)]
        table_header(cols)
        widths = [c[1] for c in cols]
        for i, p in enumerate(numeric_profiles[:15]):
            table_row([
                p.get("column", ""),
                _safe(p.get("mean")),
                _safe(p.get("std")),
                _safe(p.get("min")),
                _safe(p.get("median")),
                _safe(p.get("max")),
                _safe(p.get("missing_pct")),
            ], widths, alt=i % 2 == 1)
        if len(numeric_profiles) > 15:
            pdf.set_font("Helvetica", "I", 7)
            pdf.cell(0, 5, f"  ... and {len(numeric_profiles) - 15} more numeric features", ln=True)

    pdf.ln(2)

    if cat_profiles:
        pdf.set_font("Helvetica", "B", 9)
        pdf.cell(0, 6, f"Categorical Features ({len(cat_profiles)})", ln=True)
        cols = [("Column", 45), ("Unique", 25), ("Missing%", 25), ("Top Value", 60)]
        table_header(cols)
        widths = [c[1] for c in cols]
        for i, p in enumerate(cat_profiles[:10]):
            top_vals = p.get("top_values", {})
            top_val = next(iter(top_vals), "N/A") if top_vals else "N/A"
            table_row([
                p.get("column", ""),
                _safe(p.get("unique_count")),
                _safe(p.get("missing_pct")),
                str(top_val),
            ], widths, alt=i % 2 == 1)

    # ── Page 2 ────────────────────────────────────────────────────────────
    pdf.add_page()
    add_header_bar()

    # Correlations
    if correlations:
        section_title("Top Correlations")
        cols = [("Feature 1", 55), ("Feature 2", 55), ("Correlation", 35)]
        table_header(cols)
        widths = [c[1] for c in cols]
        sorted_corr = sorted(correlations, key=lambda x: abs(x.get("correlation", 0)), reverse=True)
        for i, c in enumerate(sorted_corr[:15]):
            table_row([
                c.get("col1", ""),
                c.get("col2", ""),
                _safe(c.get("correlation")),
            ], widths, alt=i % 2 == 1)
        pdf.ln(2)

    # Outlier analysis
    if outliers:
        section_title("Outlier Analysis")
        cols = [("Column", 50), ("IQR Outliers", 30), ("IQR %", 25), ("Z-Score Outliers", 35), ("Z-Score %", 25)]
        table_header(cols)
        widths = [c[1] for c in cols]
        sorted_outliers = sorted(outliers, key=lambda x: x.get("iqr_outlier_pct", 0), reverse=True)
        for i, o in enumerate(sorted_outliers[:12]):
            table_row([
                o.get("column", ""),
                _safe(o.get("iqr_outliers")),
                _safe(o.get("iqr_outlier_pct")),
                _safe(o.get("zscore_outliers")),
                _safe(o.get("zscore_outlier_pct")),
            ], widths, alt=i % 2 == 1)
        pdf.ln(2)

    # Feature engineering suggestions
    if feature_eng:
        section_title("Feature Engineering Suggestions")
        cols = [("Type", 35), ("Description", 120)]
        table_header(cols)
        widths = [c[1] for c in cols]
        for i, fe in enumerate(feature_eng):
            desc = fe.get("description", "")
            cols_list = fe.get("columns", [])
            if cols_list:
                desc += f"  [{', '.join(cols_list[:5])}]"
            table_row([fe.get("type", ""), _trunc(desc, 130)], widths, alt=i % 2 == 1)
        pdf.ln(2)

    # Model training results (if available)
    if training_results:
        section_title("Model Training Results")
        models = training_results if isinstance(training_results, list) else training_results.get("models", [])
        if models:
            cols = [("Model", 45), ("Accuracy", 22), ("F1", 22), ("Precision", 22), ("Recall", 22), ("AUC-ROC", 22)]
            table_header(cols)
            widths = [c[1] for c in cols]
            best_f1 = max((m.get("f1", 0) or 0 for m in models), default=0)
            for i, m in enumerate(models):
                name = m.get("model_name", m.get("name", "Unknown"))
                is_best = (m.get("f1", 0) or 0) == best_f1 and best_f1 > 0
                if is_best:
                    pdf.set_font("Helvetica", "B", 8)
                    name = f"* {name} (BEST)"
                table_row([
                    name,
                    _safe(m.get("accuracy")),
                    _safe(m.get("f1")),
                    _safe(m.get("precision")),
                    _safe(m.get("recall")),
                    _safe(m.get("auc_roc", m.get("auc-roc", m.get("roc_auc")))),
                ], widths, alt=i % 2 == 1)
                pdf.set_font("Helvetica", "", 8)

    # Footer
    pdf.ln(6)
    pdf.set_draw_color(*ACCENT)
    pdf.line(10, pdf.get_y(), 200, pdf.get_y())
    pdf.ln(3)
    pdf.set_font("Helvetica", "I", 7)
    pdf.set_text_color(120, 120, 120)
    pdf.cell(0, 4, f"Banking AI/ML Platform  |  Report generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}  |  Confidential", align="C", ln=True)

    return pdf.output()


# ── Excel generation ─────────────────────────────────────────────────────────

def _build_use_case_excel(uc_dir: Path) -> bytes:
    """Generate a multi-sheet Excel workbook for a single use case."""
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

    summary = _load_json(uc_dir / "summary.json") or {}
    column_profiles = _load_json(uc_dir / "column_profiles.json") or []
    correlations = _load_json(uc_dir / "correlations.json") or []
    feature_eng = _load_json(uc_dir / "feature_engineering.json") or []
    target_dist = _load_json(uc_dir / "target_distribution.json") or {}
    training_results = _load_json(uc_dir / "training_results.json")
    full_report = _load_json(uc_dir / "full_report.json") or {}
    outliers = _load_json(uc_dir / "outliers.json") or []

    wb = Workbook()

    # Styles
    header_font = Font(bold=True, color="FFFFFF", size=11)
    header_fill = PatternFill(start_color="0F3460", end_color="0F3460", fill_type="solid")
    accent_fill = PatternFill(start_color="E6F0FA", end_color="E6F0FA", fill_type="solid")
    thin_border = Border(
        left=Side(style="thin", color="D0D0D0"),
        right=Side(style="thin", color="D0D0D0"),
        top=Side(style="thin", color="D0D0D0"),
        bottom=Side(style="thin", color="D0D0D0"),
    )

    def style_header(ws, row: int, max_col: int):
        for col in range(1, max_col + 1):
            cell = ws.cell(row=row, column=col)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal="center")
            cell.border = thin_border

    def auto_width(ws):
        from openpyxl.cell.cell import MergedCell
        from openpyxl.utils import get_column_letter
        for col_idx, col in enumerate(ws.columns, 1):
            max_length = 0
            for cell in col:
                if isinstance(cell, MergedCell):
                    continue
                try:
                    max_length = max(max_length, len(str(cell.value or "")))
                except Exception:
                    pass
            ws.column_dimensions[get_column_letter(col_idx)].width = min(max_length + 3, 45)

    # ── Sheet 1: Summary ──────────────────────────────────────────────────
    ws = wb.active
    ws.title = "Summary"

    ws.append(["Banking AI/ML Platform - Use Case Report"])
    ws.merge_cells("A1:B1")
    ws["A1"].font = Font(bold=True, size=14, color="0F3460")
    ws.append(["Generated:", datetime.now().strftime("%Y-%m-%d %H:%M:%S")])
    ws.append([])

    ws.append(["Metric", "Value"])
    style_header(ws, 4, 2)

    dq = full_report.get("data_quality", {})
    rows = [
        ("Use Case Key", summary.get("use_case_key", "")),
        ("Label", summary.get("label", "")),
        ("Category", summary.get("category", "")),
        ("Domain", summary.get("domain", "")),
        ("Run Timestamp", summary.get("run_timestamp", "")),
        ("Data Quality Score", summary.get("data_quality_score", "")),
        ("Total Rows", summary.get("total_rows", "")),
        ("Total Columns", summary.get("total_columns", "")),
        ("Numeric Columns", summary.get("numeric_columns", "")),
        ("Categorical Columns", summary.get("categorical_columns", "")),
        ("Avg Missing %", summary.get("avg_missing_pct", "")),
        ("Duplicate Row %", dq.get("duplicate_row_pct", summary.get("duplicate_row_pct", ""))),
        ("Avg Outlier %", summary.get("avg_outlier_pct", "")),
        ("Target Column", target_dist.get("target_column", summary.get("target_column", ""))),
        ("Number of Classes", target_dist.get("n_classes", "")),
        ("Imbalance Ratio", target_dist.get("imbalance_ratio", summary.get("class_imbalance_ratio", ""))),
        ("Feature Suggestions", summary.get("n_feature_suggestions", "")),
    ]
    for r in rows:
        ws.append(list(r))

    # Class distribution sub-table
    class_dist = target_dist.get("class_distribution", {})
    if class_dist:
        ws.append([])
        ws.append(["Class", "Count"])
        style_header(ws, ws.max_row, 2)
        for cls, cnt in class_dist.items():
            ws.append([cls, cnt])

    auto_width(ws)

    # ── Sheet 2: Column Profiles ──────────────────────────────────────────
    ws2 = wb.create_sheet("Column Profiles")
    col_headers = [
        "Column", "Dtype", "Inferred Type", "Total", "Missing", "Missing%",
        "Unique", "Unique%", "Mean", "Std", "Min", "Q1", "Median", "Q3",
        "Max", "Skewness", "Kurtosis", "IQR Outliers", "IQR%", "ZScore Outliers", "ZScore%",
    ]
    ws2.append(col_headers)
    style_header(ws2, 1, len(col_headers))

    for p in column_profiles:
        ws2.append([
            p.get("column", ""),
            p.get("dtype", ""),
            p.get("inferred_type", ""),
            p.get("total_count", ""),
            p.get("missing_count", ""),
            p.get("missing_pct", ""),
            p.get("unique_count", ""),
            p.get("unique_pct", ""),
            p.get("mean", ""),
            p.get("std", ""),
            p.get("min", ""),
            p.get("q1", ""),
            p.get("median", ""),
            p.get("q3", ""),
            p.get("max", ""),
            p.get("skewness", ""),
            p.get("kurtosis", ""),
            p.get("outlier_count_iqr", ""),
            p.get("outlier_pct_iqr", ""),
            p.get("outlier_count_zscore", ""),
            p.get("outlier_pct_zscore", ""),
        ])

    auto_width(ws2)

    # ── Sheet 3: Correlations ─────────────────────────────────────────────
    ws3 = wb.create_sheet("Correlations")
    ws3.append(["Feature 1", "Feature 2", "Correlation"])
    style_header(ws3, 1, 3)
    sorted_corr = sorted(correlations, key=lambda x: abs(x.get("correlation", 0)), reverse=True)
    for c in sorted_corr:
        ws3.append([c.get("col1", ""), c.get("col2", ""), c.get("correlation", "")])
    auto_width(ws3)

    # ── Sheet 4: Model Results ────────────────────────────────────────────
    ws4 = wb.create_sheet("Model Results")
    if training_results:
        models = training_results if isinstance(training_results, list) else training_results.get("models", [])
        ws4.append(["Model", "Accuracy", "F1", "Precision", "Recall", "AUC-ROC"])
        style_header(ws4, 1, 6)
        for m in models:
            ws4.append([
                m.get("model_name", m.get("name", "")),
                m.get("accuracy", ""),
                m.get("f1", ""),
                m.get("precision", ""),
                m.get("recall", ""),
                m.get("auc_roc", m.get("auc-roc", m.get("roc_auc", ""))),
            ])
    else:
        ws4.append(["No training results available for this use case."])
    auto_width(ws4)

    # ── Sheet 5: Feature Engineering ──────────────────────────────────────
    ws5 = wb.create_sheet("Feature Engineering")
    ws5.append(["Type", "Description", "Columns"])
    style_header(ws5, 1, 3)
    for fe in feature_eng:
        ws5.append([
            fe.get("type", ""),
            fe.get("description", ""),
            ", ".join(fe.get("columns", [])),
        ])
    auto_width(ws5)

    # ── Sheet 6: Outliers ─────────────────────────────────────────────────
    ws6 = wb.create_sheet("Outliers")
    ws6.append(["Column", "IQR Outliers", "IQR %", "Z-Score Outliers", "Z-Score %"])
    style_header(ws6, 1, 5)
    for o in sorted(outliers, key=lambda x: x.get("iqr_outlier_pct", 0), reverse=True):
        ws6.append([
            o.get("column", ""),
            o.get("iqr_outliers", ""),
            o.get("iqr_outlier_pct", ""),
            o.get("zscore_outliers", ""),
            o.get("zscore_outlier_pct", ""),
        ])
    auto_width(ws6)

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ── Executive Summary PDF ────────────────────────────────────────────────────

def _build_executive_summary_pdf(preprocess_dir: Path) -> bytes:
    """Cross-use-case executive summary PDF."""
    from fpdf import FPDF

    # Gather all summaries
    summaries: List[Dict[str, Any]] = []
    if preprocess_dir.exists():
        for d in sorted(preprocess_dir.iterdir()):
            if not d.is_dir():
                continue
            s = _load_json(d / "summary.json")
            if s:
                summaries.append(s)

    if not summaries:
        raise NotFoundError("No preprocessed use cases found.")

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=20)

    HEADER_BG = (15, 52, 96)
    HEADER_FG = (255, 255, 255)
    SECTION_BG = (230, 240, 250)
    TABLE_HEADER_BG = (41, 82, 130)
    TABLE_HEADER_FG = (255, 255, 255)
    ROW_ALT = (245, 248, 252)
    ACCENT = (22, 119, 255)
    GREEN = (16, 150, 72)
    RED = (200, 40, 40)

    def add_header_bar():
        pdf.set_fill_color(*HEADER_BG)
        pdf.rect(0, 0, 210, 32, "F")
        pdf.set_xy(10, 6)
        pdf.set_font("Helvetica", "B", 16)
        pdf.set_text_color(*HEADER_FG)
        pdf.cell(0, 9, "Banking AI/ML Platform", ln=True)
        pdf.set_font("Helvetica", "", 11)
        pdf.set_x(10)
        pdf.cell(0, 6, "Executive Summary Report", ln=True)
        pdf.set_font("Helvetica", "", 8)
        pdf.set_x(10)
        pdf.cell(0, 5, f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}  |  {len(summaries)} Use Cases Analyzed", ln=True)
        pdf.set_text_color(0, 0, 0)
        pdf.set_y(38)

    def section_title(title: str):
        pdf.ln(3)
        pdf.set_fill_color(*SECTION_BG)
        pdf.set_font("Helvetica", "B", 12)
        pdf.set_text_color(*HEADER_BG)
        pdf.cell(0, 9, f"  {title}", ln=True, fill=True)
        pdf.set_text_color(0, 0, 0)
        pdf.ln(2)

    def table_header(cols: List[tuple]):
        pdf.set_fill_color(*TABLE_HEADER_BG)
        pdf.set_text_color(*TABLE_HEADER_FG)
        pdf.set_font("Helvetica", "B", 8)
        for label, w in cols:
            pdf.cell(w, 7, label, border=0, fill=True)
        pdf.ln()
        pdf.set_text_color(0, 0, 0)

    def table_row(values: List[str], widths: List[int], alt: bool = False):
        if alt:
            pdf.set_fill_color(*ROW_ALT)
        else:
            pdf.set_fill_color(255, 255, 255)
        pdf.set_font("Helvetica", "", 7)
        for val, w in zip(values, widths):
            pdf.cell(w, 6, _trunc(val, int(w * 1.2)), border=0, fill=True)
        pdf.ln()

    # ── Page 1 ────────────────────────────────────────────────────────────
    pdf.add_page()
    add_header_bar()

    # Key metrics
    section_title("Platform Overview")
    total_rows = sum(s.get("total_rows", 0) for s in summaries)
    avg_quality = sum(s.get("data_quality_score", 0) for s in summaries) / len(summaries)
    scores = [s.get("data_quality_score", 0) for s in summaries]
    high_q = sum(1 for sc in scores if sc >= 85)
    med_q = sum(1 for sc in scores if 70 <= sc < 85)
    low_q = sum(1 for sc in scores if sc < 70)

    pdf.set_font("Helvetica", "", 9)

    # Metrics in a grid-like format
    metrics = [
        ("Total Use Cases", str(len(summaries))),
        ("Total Data Rows", f"{total_rows:,}"),
        ("Avg Quality Score", f"{avg_quality:.1f} / 100"),
        ("High Quality (>=85)", str(high_q)),
        ("Medium Quality (70-84)", str(med_q)),
        ("Low Quality (<70)", str(low_q)),
    ]
    for i in range(0, len(metrics), 2):
        pdf.set_font("Helvetica", "B", 9)
        pdf.cell(50, 6, metrics[i][0] + ":", border=0)
        pdf.set_font("Helvetica", "", 9)
        pdf.cell(40, 6, metrics[i][1], border=0)
        if i + 1 < len(metrics):
            pdf.set_font("Helvetica", "B", 9)
            pdf.cell(50, 6, metrics[i + 1][0] + ":", border=0)
            pdf.set_font("Helvetica", "", 9)
            pdf.cell(40, 6, metrics[i + 1][1], border=0)
        pdf.ln()

    pdf.ln(2)

    # Department breakdown
    section_title("Department Breakdown")
    dept_counts: Dict[str, list] = {}
    for s in summaries:
        cat = s.get("category", "Unknown")
        dept_counts.setdefault(cat, []).append(s)

    cols = [("Department / Category", 60), ("Use Cases", 25), ("Avg Quality", 30), ("Total Rows", 35)]
    table_header(cols)
    widths = [c[1] for c in cols]
    for i, (dept, items) in enumerate(sorted(dept_counts.items())):
        avg_q = sum(x.get("data_quality_score", 0) for x in items) / len(items)
        tot_r = sum(x.get("total_rows", 0) for x in items)
        table_row([dept, str(len(items)), f"{avg_q:.1f}", f"{tot_r:,}"], widths, alt=i % 2 == 1)

    pdf.ln(2)

    # Top 10 by quality
    section_title("Top 10 Use Cases by Data Quality Score")
    sorted_by_quality = sorted(summaries, key=lambda x: x.get("data_quality_score", 0), reverse=True)[:10]
    cols = [("Rank", 12), ("Use Case", 70), ("Category", 40), ("Score", 20), ("Rows", 25)]
    table_header(cols)
    widths = [c[1] for c in cols]
    for i, s in enumerate(sorted_by_quality):
        table_row([
            str(i + 1),
            s.get("label", s.get("use_case_key", "")),
            s.get("category", ""),
            _safe(s.get("data_quality_score")),
            f"{s.get('total_rows', 0):,}",
        ], widths, alt=i % 2 == 1)

    # ── Page 2 ────────────────────────────────────────────────────────────
    pdf.add_page()
    add_header_bar()

    # Top 10 by data size
    section_title("Top 10 Use Cases by Data Size")
    sorted_by_size = sorted(summaries, key=lambda x: x.get("total_rows", 0), reverse=True)[:10]
    cols = [("Rank", 12), ("Use Case", 70), ("Category", 40), ("Rows", 28), ("Columns", 20)]
    table_header(cols)
    widths = [c[1] for c in cols]
    for i, s in enumerate(sorted_by_size):
        table_row([
            str(i + 1),
            s.get("label", s.get("use_case_key", "")),
            s.get("category", ""),
            f"{s.get('total_rows', 0):,}",
            _safe(s.get("total_columns")),
        ], widths, alt=i % 2 == 1)

    pdf.ln(3)

    # Data quality distribution
    section_title("Data Quality Distribution")
    quality_bands = [
        ("Excellent (95-100)", 95, 100),
        ("Good (85-94)", 85, 95),
        ("Acceptable (70-84)", 70, 85),
        ("Poor (50-69)", 50, 70),
        ("Critical (<50)", 0, 50),
    ]
    cols = [("Quality Band", 55), ("Count", 20), ("Percentage", 30)]
    table_header(cols)
    widths = [c[1] for c in cols]
    for i, (band, lo, hi) in enumerate(quality_bands):
        count = sum(1 for s in summaries if lo <= s.get("data_quality_score", 0) < hi)
        pct = count / len(summaries) * 100 if summaries else 0
        table_row([band, str(count), f"{pct:.1f}%"], widths, alt=i % 2 == 1)

    pdf.ln(3)

    # Imbalance overview
    section_title("Target Imbalance Overview")
    imbalanced = [(s, s.get("class_imbalance_ratio", 0)) for s in summaries if s.get("class_imbalance_ratio")]
    imbalanced.sort(key=lambda x: x[1], reverse=True)
    if imbalanced:
        cols = [("Use Case", 75), ("Target", 35), ("Imbalance Ratio", 30)]
        table_header(cols)
        widths = [c[1] for c in cols]
        for i, (s, ratio) in enumerate(imbalanced[:10]):
            table_row([
                s.get("label", s.get("use_case_key", "")),
                s.get("target_column", "N/A"),
                f"{ratio:.1f}:1",
            ], widths, alt=i % 2 == 1)
    else:
        pdf.set_font("Helvetica", "I", 9)
        pdf.cell(0, 6, "No class imbalance data available.", ln=True)

    # Footer
    pdf.ln(6)
    pdf.set_draw_color(*ACCENT)
    pdf.line(10, pdf.get_y(), 200, pdf.get_y())
    pdf.ln(3)
    pdf.set_font("Helvetica", "I", 7)
    pdf.set_text_color(120, 120, 120)
    pdf.cell(0, 4, f"Banking AI/ML Platform  |  Executive Summary  |  {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}  |  Confidential", align="C", ln=True)

    return pdf.output()


# ── Pydantic models ──────────────────────────────────────────────────────────

class BatchExportRequest(BaseModel):
    uc_ids: List[str]
    format: str = "pdf"  # "pdf" or "excel"


# ── Endpoints ────────────────────────────────────────────────────────────────

@router.post("/pdf/{uc_id}")
async def export_pdf(uc_id: str, settings: Settings = Depends(get_settings)):
    """Generate and download a PDF report for a single use case."""
    uc_dir = _find_uc_dir(uc_id, settings.output_dir)
    if not uc_dir:
        raise NotFoundError(f"Use case '{uc_id}' not found in preprocessing output.")

    try:
        pdf_bytes = _build_use_case_pdf(uc_dir)
    except Exception as e:
        logger.exception("PDF generation failed for %s", uc_id)
        raise DataError(f"PDF generation failed: {e}")

    filename = f"{uc_dir.name}_report.pdf"
    return StreamingResponse(
        io.BytesIO(pdf_bytes),
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@router.post("/excel/{uc_id}")
async def export_excel(uc_id: str, settings: Settings = Depends(get_settings)):
    """Generate and download an Excel workbook for a single use case."""
    uc_dir = _find_uc_dir(uc_id, settings.output_dir)
    if not uc_dir:
        raise NotFoundError(f"Use case '{uc_id}' not found in preprocessing output.")

    try:
        xlsx_bytes = _build_use_case_excel(uc_dir)
    except Exception as e:
        logger.exception("Excel generation failed for %s", uc_id)
        raise DataError(f"Excel generation failed: {e}")

    filename = f"{uc_dir.name}_report.xlsx"
    return StreamingResponse(
        io.BytesIO(xlsx_bytes),
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@router.post("/executive-summary")
async def export_executive_summary(settings: Settings = Depends(get_settings)):
    """Generate a cross-use-case executive summary PDF."""
    try:
        pdf_bytes = _build_executive_summary_pdf(settings.output_dir)
    except (NotFoundError, ValidationError, DataError):
        raise
    except Exception as e:
        logger.exception("Executive summary generation failed")
        raise DataError(f"Executive summary generation failed: {e}")

    filename = f"executive_summary_{datetime.now().strftime('%Y%m%d_%H%M')}.pdf"
    return StreamingResponse(
        io.BytesIO(pdf_bytes),
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@router.post("/batch")
async def export_batch(request: BatchExportRequest, settings: Settings = Depends(get_settings)):
    """Generate a ZIP archive containing individual reports for multiple use cases."""
    if not request.uc_ids:
        raise ValidationError("uc_ids list cannot be empty.")
    if request.format not in ("pdf", "excel"):
        raise ValidationError("format must be 'pdf' or 'excel'.")

    buf = io.BytesIO()
    errors: List[str] = []

    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for uc_id in request.uc_ids:
            uc_dir = _find_uc_dir(uc_id, settings.output_dir)
            if not uc_dir:
                errors.append(f"{uc_id}: not found")
                continue
            try:
                if request.format == "pdf":
                    data = _build_use_case_pdf(uc_dir)
                    ext = "pdf"
                else:
                    data = _build_use_case_excel(uc_dir)
                    ext = "xlsx"
                zf.writestr(f"{uc_dir.name}_report.{ext}", data)
            except Exception as e:
                errors.append(f"{uc_id}: {e}")
                logger.warning("Batch export error for %s: %s", uc_id, e)

        # Add a manifest with any errors
        if errors:
            manifest = "Export Errors:\n" + "\n".join(errors)
            zf.writestr("_export_errors.txt", manifest)

    buf.seek(0)
    fmt_label = "pdf" if request.format == "pdf" else "xlsx"
    filename = f"batch_export_{fmt_label}_{datetime.now().strftime('%Y%m%d_%H%M')}.zip"
    return StreamingResponse(
        buf,
        media_type="application/zip",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


# ── Word (.docx) generation ──────────────────────────────────────────────────

def _build_use_case_word(uc_dir: Path) -> bytes:
    """Generate a Word document report for a single use case."""
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.table import WD_TABLE_ALIGNMENT

    summary = _load_json(uc_dir / "summary.json") or {}
    full_report = _load_json(uc_dir / "full_report.json") or {}
    column_profiles = _load_json(uc_dir / "column_profiles.json") or []
    feature_eng = _load_json(uc_dir / "feature_engineering.json") or []
    outliers = _load_json(uc_dir / "outliers.json") or []
    target_dist = _load_json(uc_dir / "target_distribution.json") or {}
    correlations = _load_json(uc_dir / "correlations.json") or []
    training_results = _load_json(uc_dir / "training_results.json")
    dq = full_report.get("data_quality", {})
    label = summary.get("label", uc_dir.name)

    doc = Document()

    # Title
    title = doc.add_heading("Banking AI/ML Platform — Use Case Report", level=0)
    title.runs[0].font.color.rgb = RGBColor(15, 52, 96)
    doc.add_paragraph(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}")

    # Use Case Overview
    doc.add_heading("Use Case Overview", level=1)
    for k, v in [
        ("Name", label),
        ("Key", summary.get("use_case_key", "N/A")),
        ("Category", summary.get("category", "N/A")),
        ("Domain", summary.get("domain", "N/A")),
        ("Run Timestamp", summary.get("run_timestamp", "N/A")),
    ]:
        p = doc.add_paragraph()
        p.add_run(f"{k}: ").bold = True
        p.add_run(_safe(v))

    # Data Quality
    doc.add_heading("Data Quality", level=1)
    tbl = doc.add_table(rows=1, cols=2)
    tbl.style = "Light Grid Accent 1"
    tbl.alignment = WD_TABLE_ALIGNMENT.LEFT
    hdr = tbl.rows[0].cells
    hdr[0].text, hdr[1].text = "Metric", "Value"
    for k, v in [
        ("Total Rows", _safe(summary.get("total_rows"))),
        ("Total Columns", _safe(summary.get("total_columns"))),
        ("Numeric Columns", _safe(summary.get("numeric_columns"))),
        ("Categorical Columns", _safe(summary.get("categorical_columns"))),
        ("Quality Score", f"{_safe(summary.get('data_quality_score'))} / 100"),
        ("Avg Missing %", _safe(summary.get("avg_missing_pct"))),
        ("Duplicate Row %", _safe(dq.get("duplicate_row_pct", summary.get("duplicate_row_pct")))),
        ("Avg Outlier %", _safe(summary.get("avg_outlier_pct"))),
    ]:
        row = tbl.add_row().cells
        row[0].text, row[1].text = k, v

    # Target Analysis
    doc.add_heading("Target Analysis", level=1)
    for k, v in [
        ("Target Column", _safe(target_dist.get("target_column", summary.get("target_column")))),
        ("Number of Classes", _safe(target_dist.get("n_classes"))),
        ("Imbalance Ratio", _safe(target_dist.get("imbalance_ratio", summary.get("class_imbalance_ratio")))),
    ]:
        p = doc.add_paragraph()
        p.add_run(f"{k}: ").bold = True
        p.add_run(v)

    class_dist = target_dist.get("class_distribution", {})
    if class_dist:
        tbl2 = doc.add_table(rows=1, cols=3)
        tbl2.style = "Light Grid Accent 1"
        h = tbl2.rows[0].cells
        h[0].text, h[1].text, h[2].text = "Class", "Count", "Proportion"
        total = sum(class_dist.values()) or 1
        for cls, cnt in class_dist.items():
            r = tbl2.add_row().cells
            r[0].text, r[1].text, r[2].text = str(cls), str(cnt), f"{cnt/total*100:.1f}%"

    # Feature Summary
    numeric_profiles = [p for p in column_profiles if p.get("inferred_type") == "numeric"]
    if numeric_profiles:
        doc.add_heading("Numeric Features", level=2)
        tbl3 = doc.add_table(rows=1, cols=5)
        tbl3.style = "Light Grid Accent 1"
        for i, h in enumerate(["Column", "Mean", "Std", "Min", "Max"]):
            tbl3.rows[0].cells[i].text = h
        for p in numeric_profiles[:20]:
            r = tbl3.add_row().cells
            r[0].text = p.get("column", "")
            r[1].text, r[2].text = _safe(p.get("mean")), _safe(p.get("std"))
            r[3].text, r[4].text = _safe(p.get("min")), _safe(p.get("max"))

    # Correlations
    if correlations:
        doc.add_heading("Top Correlations", level=1)
        sorted_corr = sorted(correlations, key=lambda x: abs(x.get("correlation", 0)), reverse=True)[:15]
        tbl4 = doc.add_table(rows=1, cols=3)
        tbl4.style = "Light Grid Accent 1"
        for i, h in enumerate(["Feature 1", "Feature 2", "Correlation"]):
            tbl4.rows[0].cells[i].text = h
        for c in sorted_corr:
            r = tbl4.add_row().cells
            r[0].text, r[1].text, r[2].text = c.get("col1", ""), c.get("col2", ""), _safe(c.get("correlation"))

    # Model Results
    if training_results:
        doc.add_heading("Model Training Results", level=1)
        models = training_results if isinstance(training_results, list) else training_results.get("models", [])
        if models:
            tbl5 = doc.add_table(rows=1, cols=6)
            tbl5.style = "Light Grid Accent 1"
            for i, h in enumerate(["Model", "Accuracy", "F1", "Precision", "Recall", "AUC-ROC"]):
                tbl5.rows[0].cells[i].text = h
            for m in models:
                r = tbl5.add_row().cells
                r[0].text = m.get("model_name", m.get("name", ""))
                r[1].text, r[2].text = _safe(m.get("accuracy")), _safe(m.get("f1"))
                r[3].text, r[4].text = _safe(m.get("precision")), _safe(m.get("recall"))
                r[5].text = _safe(m.get("auc_roc", m.get("auc-roc", m.get("roc_auc"))))

    doc.add_paragraph(f"\nGenerated by Banking AI/ML Platform on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ── Markdown generation ───────────────────────────────────────────────────────

def _build_use_case_markdown(uc_dir: Path) -> str:
    """Generate a Markdown report for a single use case."""
    summary = _load_json(uc_dir / "summary.json") or {}
    full_report = _load_json(uc_dir / "full_report.json") or {}
    column_profiles = _load_json(uc_dir / "column_profiles.json") or []
    target_dist = _load_json(uc_dir / "target_distribution.json") or {}
    correlations = _load_json(uc_dir / "correlations.json") or []
    training_results = _load_json(uc_dir / "training_results.json")
    dq = full_report.get("data_quality", {})
    label = summary.get("label", uc_dir.name)

    lines = [
        f"# Banking AI/ML Platform — Use Case Report",
        f"",
        f"**Generated:** {datetime.now().strftime('%Y-%m-%d %H:%M')}",
        f"",
        f"## Use Case Overview",
        f"",
        f"| Field | Value |",
        f"|-------|-------|",
        f"| Name | {label} |",
        f"| Key | {summary.get('use_case_key', 'N/A')} |",
        f"| Category | {summary.get('category', 'N/A')} |",
        f"| Domain | {summary.get('domain', 'N/A')} |",
        f"| Run Timestamp | {summary.get('run_timestamp', 'N/A')} |",
        f"",
        f"## Data Quality",
        f"",
        f"| Metric | Value |",
        f"|--------|-------|",
        f"| Total Rows | {_safe(summary.get('total_rows'))} |",
        f"| Total Columns | {_safe(summary.get('total_columns'))} |",
        f"| Numeric Columns | {_safe(summary.get('numeric_columns'))} |",
        f"| Categorical Columns | {_safe(summary.get('categorical_columns'))} |",
        f"| Quality Score | {_safe(summary.get('data_quality_score'))} / 100 |",
        f"| Avg Missing % | {_safe(summary.get('avg_missing_pct'))} |",
        f"| Duplicate Row % | {_safe(dq.get('duplicate_row_pct', summary.get('duplicate_row_pct')))} |",
        f"| Avg Outlier % | {_safe(summary.get('avg_outlier_pct'))} |",
        f"",
        f"## Target Analysis",
        f"",
        f"- **Target Column:** {_safe(target_dist.get('target_column', summary.get('target_column')))}",
        f"- **Number of Classes:** {_safe(target_dist.get('n_classes'))}",
        f"- **Imbalance Ratio:** {_safe(target_dist.get('imbalance_ratio', summary.get('class_imbalance_ratio')))}",
        f"",
    ]

    class_dist = target_dist.get("class_distribution", {})
    if class_dist:
        lines += ["| Class | Count | Proportion |", "|-------|-------|------------|"]
        total = sum(class_dist.values()) or 1
        for cls, cnt in class_dist.items():
            lines.append(f"| {cls} | {cnt} | {cnt/total*100:.1f}% |")
        lines.append("")

    # Numeric features
    numeric_profiles = [p for p in column_profiles if p.get("inferred_type") == "numeric"]
    if numeric_profiles:
        lines += [f"## Numeric Features ({len(numeric_profiles)})", "",
                  "| Column | Mean | Std | Min | Max |",
                  "|--------|------|-----|-----|-----|"]
        for p in numeric_profiles[:20]:
            lines.append(f"| {p.get('column','')} | {_safe(p.get('mean'))} | {_safe(p.get('std'))} | {_safe(p.get('min'))} | {_safe(p.get('max'))} |")
        lines.append("")

    # Correlations
    if correlations:
        sorted_corr = sorted(correlations, key=lambda x: abs(x.get("correlation", 0)), reverse=True)[:15]
        lines += ["## Top Correlations", "",
                  "| Feature 1 | Feature 2 | Correlation |",
                  "|-----------|-----------|-------------|"]
        for c in sorted_corr:
            lines.append(f"| {c.get('col1','')} | {c.get('col2','')} | {_safe(c.get('correlation'))} |")
        lines.append("")

    # Model results
    if training_results:
        models = training_results if isinstance(training_results, list) else training_results.get("models", [])
        if models:
            lines += ["## Model Training Results", "",
                      "| Model | Accuracy | F1 | Precision | Recall | AUC-ROC |",
                      "|-------|----------|----|-----------|---------| --------|"]
            for m in models:
                name = m.get("model_name", m.get("name", ""))
                lines.append(f"| {name} | {_safe(m.get('accuracy'))} | {_safe(m.get('f1'))} | {_safe(m.get('precision'))} | {_safe(m.get('recall'))} | {_safe(m.get('auc_roc', m.get('auc-roc', m.get('roc_auc'))))} |")
            lines.append("")

    lines += [f"---", f"*Generated by Banking AI/ML Platform on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}*"]
    return "\n".join(lines)


# ── PowerPoint generation ─────────────────────────────────────────────────────

def _build_use_case_pptx(uc_dir: Path) -> bytes:
    """Generate a PowerPoint presentation for a single use case."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor as PptxRGB
    from pptx.enum.text import PP_ALIGN

    summary = _load_json(uc_dir / "summary.json") or {}
    full_report = _load_json(uc_dir / "full_report.json") or {}
    column_profiles = _load_json(uc_dir / "column_profiles.json") or []
    target_dist = _load_json(uc_dir / "target_distribution.json") or {}
    correlations = _load_json(uc_dir / "correlations.json") or []
    training_results = _load_json(uc_dir / "training_results.json")
    dq = full_report.get("data_quality", {})
    label = summary.get("label", uc_dir.name)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    NAVY = PptxRGB(15, 52, 96)
    WHITE = PptxRGB(255, 255, 255)
    ACCENT = PptxRGB(22, 119, 255)

    def add_title_slide(title, subtitle):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        # Background
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = NAVY
        # Title
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT
        # Subtitle
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = PptxRGB(180, 200, 230)
        p2.alignment = PP_ALIGN.LEFT
        return slide

    def add_content_slide(title, items):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Title bar
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1))  # rectangle
        shape.fill.solid()
        shape.fill.fore_color.rgb = NAVY
        shape.line.fill.background()
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.7))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.color.rgb = WHITE
        p.font.bold = True
        # Content
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for i, (k, v) in enumerate(items):
            p = tf2.add_paragraph() if i > 0 else tf2.paragraphs[0]
            run_k = p.add_run()
            run_k.text = f"{k}: "
            run_k.font.bold = True
            run_k.font.size = Pt(16)
            run_k.font.color.rgb = NAVY
            run_v = p.add_run()
            run_v.text = str(v)
            run_v.font.size = Pt(16)
            p.space_after = Pt(6)
        return slide

    def add_table_slide(title, headers, rows):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = NAVY
        shape.line.fill.background()
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.7))
        p = txBox.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.color.rgb = WHITE
        p.font.bold = True

        n_rows = min(len(rows) + 1, 16)
        n_cols = len(headers)
        tbl_shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.5), Inches(1.3), Inches(12), Inches(5.5))
        tbl = tbl_shape.table
        for ci, h in enumerate(headers):
            cell = tbl.cell(0, ci)
            cell.text = h
            for para in cell.text_frame.paragraphs:
                para.font.bold = True
                para.font.size = Pt(12)
                para.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
        for ri, row in enumerate(rows[:n_rows - 1]):
            for ci, val in enumerate(row):
                cell = tbl.cell(ri + 1, ci)
                cell.text = str(val)
                for para in cell.text_frame.paragraphs:
                    para.font.size = Pt(11)
        return slide

    # Slide 1: Title
    add_title_slide(
        f"Use Case Report: {label}",
        f"Banking AI/ML Platform  |  {datetime.now().strftime('%Y-%m-%d %H:%M')}"
    )

    # Slide 2: Overview
    add_content_slide("Use Case Overview", [
        ("Name", label),
        ("Key", summary.get("use_case_key", "N/A")),
        ("Category", summary.get("category", "N/A")),
        ("Domain", summary.get("domain", "N/A")),
        ("Run Timestamp", summary.get("run_timestamp", "N/A")),
    ])

    # Slide 3: Data Quality
    add_content_slide("Data Quality", [
        ("Total Rows", _safe(summary.get("total_rows"))),
        ("Total Columns", _safe(summary.get("total_columns"))),
        ("Numeric / Categorical", f"{_safe(summary.get('numeric_columns'))} / {_safe(summary.get('categorical_columns'))}"),
        ("Quality Score", f"{_safe(summary.get('data_quality_score'))} / 100"),
        ("Avg Missing %", _safe(summary.get("avg_missing_pct"))),
        ("Duplicate Row %", _safe(dq.get("duplicate_row_pct", summary.get("duplicate_row_pct")))),
        ("Avg Outlier %", _safe(summary.get("avg_outlier_pct"))),
    ])

    # Slide 4: Target Analysis
    target_items = [
        ("Target Column", _safe(target_dist.get("target_column", summary.get("target_column")))),
        ("Number of Classes", _safe(target_dist.get("n_classes"))),
        ("Imbalance Ratio", _safe(target_dist.get("imbalance_ratio", summary.get("class_imbalance_ratio")))),
    ]
    class_dist = target_dist.get("class_distribution", {})
    if class_dist:
        total = sum(class_dist.values()) or 1
        for cls, cnt in list(class_dist.items())[:8]:
            target_items.append((f"Class '{cls}'", f"{cnt} ({cnt/total*100:.1f}%)"))
    add_content_slide("Target Analysis", target_items)

    # Slide 5: Correlations table
    if correlations:
        sorted_corr = sorted(correlations, key=lambda x: abs(x.get("correlation", 0)), reverse=True)[:14]
        add_table_slide("Top Correlations",
            ["Feature 1", "Feature 2", "Correlation"],
            [[c.get("col1", ""), c.get("col2", ""), _safe(c.get("correlation"))] for c in sorted_corr])

    # Slide 6: Model Results
    if training_results:
        models = training_results if isinstance(training_results, list) else training_results.get("models", [])
        if models:
            add_table_slide("Model Training Results",
                ["Model", "Accuracy", "F1", "Precision", "Recall", "AUC-ROC"],
                [[m.get("model_name", m.get("name", "")),
                  _safe(m.get("accuracy")), _safe(m.get("f1")),
                  _safe(m.get("precision")), _safe(m.get("recall")),
                  _safe(m.get("auc_roc", m.get("auc-roc", m.get("roc_auc"))))]
                 for m in models])

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── Word endpoint ─────────────────────────────────────────────────────────────

@router.post("/word/{uc_id}")
async def export_word(uc_id: str, settings: Settings = Depends(get_settings)):
    """Generate and download a Word document report for a single use case."""
    uc_dir = _find_uc_dir(uc_id, settings.output_dir)
    if not uc_dir:
        raise NotFoundError(f"Use case '{uc_id}' not found.")
    try:
        docx_bytes = _build_use_case_word(uc_dir)
    except Exception as e:
        logger.exception("Word generation failed for %s", uc_id)
        raise DataError(f"Word generation failed: {e}")
    filename = f"{uc_dir.name}_report.docx"
    return StreamingResponse(
        io.BytesIO(docx_bytes),
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


# ── Markdown endpoint ─────────────────────────────────────────────────────────

@router.post("/markdown/{uc_id}")
async def export_markdown(uc_id: str, settings: Settings = Depends(get_settings)):
    """Generate and download a Markdown report for a single use case."""
    uc_dir = _find_uc_dir(uc_id, settings.output_dir)
    if not uc_dir:
        raise NotFoundError(f"Use case '{uc_id}' not found.")
    try:
        md_text = _build_use_case_markdown(uc_dir)
    except Exception as e:
        logger.exception("Markdown generation failed for %s", uc_id)
        raise DataError(f"Markdown generation failed: {e}")
    filename = f"{uc_dir.name}_report.md"
    return StreamingResponse(
        io.BytesIO(md_text.encode("utf-8")),
        media_type="text/markdown",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


# ── PowerPoint endpoint ──────────────────────────────────────────────────────

@router.post("/pptx/{uc_id}")
async def export_pptx(uc_id: str, settings: Settings = Depends(get_settings)):
    """Generate and download a PowerPoint presentation for a single use case."""
    uc_dir = _find_uc_dir(uc_id, settings.output_dir)
    if not uc_dir:
        raise NotFoundError(f"Use case '{uc_id}' not found.")
    try:
        pptx_bytes = _build_use_case_pptx(uc_dir)
    except Exception as e:
        logger.exception("PPTX generation failed for %s", uc_id)
        raise DataError(f"PPTX generation failed: {e}")
    filename = f"{uc_dir.name}_report.pptx"
    return StreamingResponse(
        io.BytesIO(pptx_bytes),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )
